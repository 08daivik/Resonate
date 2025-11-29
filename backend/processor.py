import assemblyai as aai
from PyPDF2 import PdfReader
from pptx import Presentation

# -------------------------------
# AUDIO TRANSCRIPTION
# -------------------------------
def transcribe_audio(path: str, api_key: str):
    aai.settings.api_key = api_key
    transcriber = aai.Transcriber()
    result = transcriber.transcribe(path)
    return result.text


# -------------------------------
# PDF EXTRACTION
# -------------------------------
def extract_pdf(path: str):
    reader = PdfReader(path)
    text = ""

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"

    return text


# -------------------------------
# PPT EXTRACTION
# -------------------------------
def extract_ppt(path: str):
    prs = Presentation(path)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text:
                    text += shape.text + "\n"

    return text
